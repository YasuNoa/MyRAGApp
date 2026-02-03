
import os
import io
import uuid
import logging
import json
import re
from typing import List, Dict, Any, Optional
from datetime import datetime, timezone, timedelta

import google.generativeai as genai
from pypdf import PdfReader
from pptx import Presentation
from docx import Document as DocxDocument
import pandas as pd
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type


from database.db import db
from services.vector_service import VectorService
from services.user_service import UserService
from services.prompts import (
    PDF_TRANSCRIPTION_PROMPT,
    IMAGE_DESCRIPTION_PROMPT
)

# Setup Logger
logger = logging.getLogger(__name__)

# Timezone Definition
JST = timezone(timedelta(hours=9))

# Initialize Clients
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

if GOOGLE_API_KEY:
    genai.configure(api_key=GOOGLE_API_KEY)

class KnowledgeService:

    @staticmethod
    def extract_text_from_pdf(file_content: bytes) -> str:
        # PDFファイルのバイナリデータからテキストを抽出します。
        # pypdfライブラリを使用して、ページごとにテキストを読み取ります。
        try:
            reader = PdfReader(io.BytesIO(file_content))
            text = ""
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {e}")
            return "" 

    @staticmethod
    async def _process_pdf_with_gemini(content: bytes) -> str:
        # PDFをGemini 2.0 Flashにアップロードして、テキスト抽出 (OCR) を行います。
        # 通常のテキスト抽出が失敗した場合や、画像中心のPDFの場合に使用します。
        temp_filename = f"/tmp/{uuid.uuid4()}.pdf"
        try:
            with open(temp_filename, "wb") as f:
                f.write(content)
            
            logger.info("Uploading PDF to Gemini for OCR...")
            uploaded_file = genai.upload_file(temp_filename, mime_type="application/pdf")
            
            logger.info("Generating PDF transcript...")
            model = genai.GenerativeModel('gemini-2.0-flash')
            prompt = PDF_TRANSCRIPTION_PROMPT
            response = model.generate_content([prompt, uploaded_file])
            
            # Check if we have a valid response
            if not response.candidates:
                logger.warning(f"Gemini PDF processing returned no candidates. Prompt feedback: {response.prompt_feedback}")
                return ""
                
            try:
                return response.text
            except ValueError:
                # response.text raises ValueError if the response was blocked
                logger.warning(f"Gemini PDF processing blocked. Prompt feedback: {response.prompt_feedback}")
                return ""
        except Exception as e:
            logger.error(f"Error processing PDF with Gemini: {e}")
            return ""
        finally:
            if os.path.exists(temp_filename):
                os.remove(temp_filename)

    @classmethod
    async def process_pdf(cls, content: bytes) -> str:
        text = cls.extract_text_from_pdf(content)
        if not text.strip() or len(text.strip()) < 50: 
            logger.info("PDF text extraction yielded little/no text. Attempting Gemini OCR...")
            ocr_text = await cls._process_pdf_with_gemini(content)
            if ocr_text.strip():
                text = ocr_text
        return text

    @staticmethod
    @retry(
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=2, max=10),
        retry=retry_if_exception_type(Exception)
    )
    async def process_image(content: bytes, mime_type: str, filename: str) -> str:
        # 画像をGeminiにアップロードして、その内容説明 (Description) を生成させます。
        # これにより、画像の内容もテキストとして検索可能になります。
        temp_filename = f"/tmp/{uuid.uuid4()}_{filename}"
        with open(temp_filename, "wb") as f:
            f.write(content)
        
        try:
            logger.info("Uploading image to Gemini...")
            uploaded_file = genai.upload_file(temp_filename, mime_type=mime_type)
            
            logger.info("Generating image description...")
            model = genai.GenerativeModel('gemini-2.0-flash')
            prompt = IMAGE_DESCRIPTION_PROMPT
            response = model.generate_content([prompt, uploaded_file])
            return response.text
        finally:
            if os.path.exists(temp_filename):
                os.remove(temp_filename)

    @staticmethod
    async def process_pptx(content: bytes) -> str:
        ppt = Presentation(io.BytesIO(content))
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    @staticmethod
    async def process_docx(content: bytes) -> str:
        doc = DocxDocument(io.BytesIO(content))
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    @staticmethod
    async def process_xlsx(content: bytes) -> str:
        xls = pd.read_excel(io.BytesIO(content), sheet_name=None)
        text = ""
        for sheet_name, df in xls.items():
            text += f"--- Sheet: {sheet_name} ---\n"
            text += df.to_string(index=False) + "\n\n"
        return text

    @staticmethod
    async def process_csv(content: bytes) -> str:
        return content.decode("utf-8")

    @staticmethod
    async def process_text_file(content: bytes) -> str:
         # Fallback decode
         try:
             return content.decode("utf-8")
         except UnicodeDecodeError:
             return content.decode("utf-8", errors="ignore")

    @staticmethod
    async def create_document_record(
        doc_id: str,
        user_id: str,
        title: str,
        course_id: Optional[str] = None, # course_idを追加
        source: str = "import",
        mime_type: Optional[str] = None,
        tags: List[str] = []
    ):
        try:
            # Check if exists (idempotency)
            existing = await db.document.find_unique(where={"id": doc_id})
            if existing:
                return

            await db.document.create(
                data={
                    "id": doc_id,
                    "userId": user_id,
                    "courseId": course_id,
                    "title": title,
                    "source": source,
                    "mimeType": mime_type,
                    "tags": tags,
                    # createdAt is auto-generated by Prisma (@default(now())) or we can let DB handle it.
                    # Prisma schema probably has @default(now())
                }
            )
            logger.info(f"Created Document record (Prisma): {doc_id} in Course: {course_id}")
        except Exception as e:
            logger.error(f"Error creating Document record: {e}")
            raise e

    @staticmethod
    async def save_document_content(
        doc_id: str, 
        content: str, 
        category: Optional[str] = None, 
        summary: Optional[str] = None,
        mime_type: Optional[str] = None
    ):
        # ドキュメントの全文テキスト、要約、MIMEタイプをPostgreSQLに保存します。
        # これにより、RAGで検索した際に、元の文脈全体を参照できるようになります (Long Context RAG)。
        try:
            await db.document.update(
                where={"id": doc_id},
                data={
                    "content": content,
                    "summary": summary,
                    "mimeType": mime_type
                }
            )
            logger.info(f"Saved content for document {doc_id}")
        except Exception as e:
            logger.error(f"Error saving content to DB: {e}")

    @staticmethod
    async def get_document_content(doc_id: str) -> str:
        # PostgreSQLからドキュメントの全文テキストを取得します。
        # RAGで回答を生成する際、検索でヒットした断片だけでなく、この全文をAIに渡すことで精度を高めます。
        try:
            doc = await db.document.find_unique(where={"id": doc_id})
            if doc and doc.content:
                return doc.content
            return ""
        except Exception as e:
            logger.error(f"Error fetching content from DB: {e}")
            return ""

    @classmethod
    async def process_and_save_content(
        cls,
        text: str, 
        metadata: dict, 
        summary: Optional[str] = None
    ):

        if not text.strip() and not summary:
            # raise HTTPException(status_code=400, detail="Extracted text is empty")
            # We are in Service layer, so maybe just raise Exception
            raise ValueError("Extracted text is empty")

        user_id = metadata.get("userId")
        # file_id is legacy, we use dbId for Document ID
        file_id = metadata.get("fileId") 
        tags = metadata.get("tags", [])
        file_name = metadata.get("fileName", "Unknown")
        db_id = metadata.get("dbId")
        mime_type = metadata.get("mimeType") # Get mimeType
        course_id = metadata.get("courseId") # メタデータからcourseIdを取得

from utils.validators import validate_course_access

        # Ensure we have a db_id (Document ID) and the Document record exists
        if not db_id:
            db_id = str(uuid.uuid4())
            metadata["dbId"] = db_id
        
        # Validate Course Access if course_id is provided
        if course_id:
            await validate_course_access(course_id, user_id)

        # Create Document Row if needed
        if user_id: # Only if we have a user to attach to
            await cls.create_document_record(
                doc_id=db_id,
                user_id=user_id,
                title=file_name,
                course_id=course_id, # serviceへ渡す
                source=metadata.get("source", "import"),
                mime_type=mime_type,
                tags=tags
            )

        # テキストのチャンク分割
        chunks = VectorService.chunk_text(text)
        logger.info(f"Generated {len(chunks)} chunks for file {file_name}")

        # ベクトル化とPgvectorへの保存
        vectors = []
        for i, chunk in enumerate(chunks):
            vector_id = f"{user_id}#{db_id}#{i}"
            embedding = VectorService.get_embedding(chunk)
            
            metadata_payload = {
                "userId": user_id,
                "fileId": db_id, # Use dbId as fileId for consistency
                "dbId": db_id,   # Explicit documentId
                "fileName": file_name,
                "text": chunk,
                "chunkIndex": i,
                "tags": tags
            }
            
            vectors.append({
                "id": vector_id,
                "values": embedding,
                "metadata": metadata_payload
            })
            
            if len(vectors) >= 100:
                await VectorService.upsert_vectors(vectors)
                vectors = []

        if vectors:
            await VectorService.upsert_vectors(vectors)

        # 全文コンテンツをDBに保存 (UPDATE)
        # create_document_recordで作成済みなので、ここではcontent, summaryを更新
        if db_id:
             # PostgreSQLはNULLバイトを許容しないため、サニタイズします
             clean_text = text.replace("\x00", "")
             clean_summary = summary.replace("\x00", "") if summary else None
             
             await cls.save_document_content(db_id, clean_text, summary=clean_summary, mime_type=mime_type)
        
        return {
            "status": "success", 
            "message": f"Successfully processed {file_name}",
            "chunks_count": len(chunks),
            "fileId": file_id
        }

    @staticmethod
    async def get_categories(user_id: str) -> List[str]:
        try:
             # Prisma doesn't support unnest easily in find_many. Use query_raw.
             # Note: query_raw returns a list of dictionaries (mostly).
             query = 'SELECT DISTINCT unnest(tags) as tag FROM "Document" WHERE "userId" = $1'
             # Prisma python query_raw params
             # await db.query_raw(query, user_id)
             # But Prisma Client Python uses ? or params list.
             # Actually cleaner:
             rows = await db.query_raw(
                 'SELECT DISTINCT unnest(tags) as tag FROM "Document" WHERE "userId" = $1', 
                 user_id
             )
             # rows is typically List[Record] or List[dict] depending on driver.
             # Let's assume list of dict-like objects.
             return [row['tag'] for row in rows if row['tag']]
        except Exception as e:
            logger.error(f"Error fetching categories: {e}")
            return []

    @staticmethod
    async def get_trash_documents(user_id: str) -> List[Dict]:
        try:
            # 30日以上前のものは自動削除するバッチが必要だが、ここでは取得時にフィルタはしない（すべて表示）
            docs = await db.document.find_many(
                where={
                    "userId": user_id,
                    "deletedAt": {"not": None}
                },
                order={"deletedAt": "desc"}
            )
            # Convert to dict for serialization
            return [doc.model_dump() for doc in docs]
        except Exception as e:
            logger.error(f"Error fetching trash documents: {e}")
            return []

    @staticmethod
    async def restore_document(doc_id: str, user_id: str):
        try:
            await db.document.update_many(
                where={"id": doc_id, "userId": user_id},
                data={"deletedAt": None}
            )
            logger.info(f"Restored Document: {doc_id}")
        except Exception as e:
            logger.error(f"Error restoring document {doc_id}: {e}")
            raise e

    @staticmethod
    async def delete_document(doc_id: str, user_id: str):
        """
        Soft Delete: Set deletedAt to NOW.
        The file remains in storage/vector DB but is hidden from normal views.
        """
        try:
            # Prisma uses timezone aware datetime for NOW typically, 
            # but usually we pass datetime.now(JST)
            # from datetime import datetime # Removed local import
            
            # Using update_many for safety (id + owner)
            await db.document.update_many(
                where={"id": doc_id, "userId": user_id},
                data={"deletedAt": datetime.now(JST)}
            )
            logger.info(f"Soft Deleted Document: {doc_id}")
        except Exception as e:
            logger.error(f"Error soft deleting document {doc_id}: {e}")
            raise e

    @staticmethod
    async def permanent_delete_document(doc_id: str, user_id: str):
        """
        Hard Delete: Remove from DB record and Vector DB.
        """
        try:
            # 1. Delete from DB (Document table)
            # Using delete_many to enforce userId check
            await db.document.delete_many(
                where={"id": doc_id, "userId": user_id}
            )
            logger.info(f"Permanently Deleted Document record: {doc_id}")

            # 2. Delete Vectors (Supabase)
            # Legacy fileId might differ, but in new logic fileId == dbId.
            # Assuming doc_id is the key used for vectors as well.
            await VectorService.delete_vectors(file_id=doc_id, user_id=user_id)
            
        except Exception as e:
            logger.error(f"Error deleting document {doc_id}: {e}")
            raise e

    @staticmethod
    async def update_knowledge(
        doc_id: str, 
        user_id: str, 
        tags: Optional[List[str]] = None, 
        title: Optional[str] = None
    ):
        """
        ドキュメントのメタデータ（タグ、タイトル）を更新します。
        """
        try:
            data = {}
            if title is not None:
                data["title"] = title
            if tags is not None:
                data["tags"] = tags
            
            if not data:
                return

            # Update DB (Using update_many for safe ownership check)
            await db.document.update_many(
                where={"id": doc_id, "userId": user_id},
                data=data
            )
            
            if tags is not None:
                # Update Vectors
                await VectorService.update_tags(file_id=doc_id, user_id=user_id, tags=tags)
                
            logger.info(f"Updated knowledge for document {doc_id}")
            
        except Exception as e:
            logger.error(f"Error updating knowledge for document {doc_id}: {e}")
            raise e
